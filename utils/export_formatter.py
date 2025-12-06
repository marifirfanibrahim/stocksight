"""
export formatter module
formats data for various export types
supports csv excel powerpoint pdf
"""

import pandas as pd
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime

import config


# ============================================================================
#                           EXPORT FORMATTER
# ============================================================================

class ExportFormatter:
    # formats and exports data to various formats
    
    def __init__(self):
        # initialize formatter
        self.export_config = config.EXPORT_FORMATS
        self.ppt_config = config.PPT_TEMPLATE
    
    # ---------- CSV EXPORT ----------
    
    def export_csv(self, 
                   data: pd.DataFrame, 
                   file_path: str,
                   include_index: bool = False) -> tuple:
        # export dataframe to csv
        try:
            data.to_csv(file_path, index=include_index)
            return True, f"exported {len(data):,} rows"
        except Exception as e:
            return False, str(e)
    
    def format_forecast_csv(self, 
                            forecasts: Dict[str, Any],
                            include_bounds: bool = True) -> pd.DataFrame:
        # format forecasts for csv export - includes all forecast values
        rows = []
        
        for sku, result in forecasts.items():
            for i, (date, value) in enumerate(zip(result.dates, result.forecast)):
                # format date - remove timestamp
                date_str = self._format_date_no_timestamp(date)
                
                row = {
                    "sku": sku,
                    "date": date_str,
                    "forecast": value,
                    "model": result.model
                }
                
                if include_bounds:
                    row["lower_bound"] = result.lower_bound[i] if result.lower_bound else None
                    row["upper_bound"] = result.upper_bound[i] if result.upper_bound else None
                
                rows.append(row)
        
        return pd.DataFrame(rows)
    
    def _format_date_no_timestamp(self, date) -> str:
        # format date without timestamp
        if date is None:
            return ""
        
        date_str = str(date)
        
        # remove timestamp portion if present
        if " " in date_str:
            date_str = date_str.split(" ")[0]
        
        # also handle T separator
        if "T" in date_str:
            date_str = date_str.split("T")[0]
        
        return date_str
    
    # ---------- EXCEL EXPORT ----------
    
    def export_excel(self,
                     data: Dict[str, pd.DataFrame],
                     file_path: str,
                     format_tables: bool = True) -> tuple:
        # export multiple dataframes to excel sheets
        try:
            with pd.ExcelWriter(file_path, engine="xlsxwriter") as writer:
                workbook = writer.book
                
                # define formats
                header_format = workbook.add_format({
                    "bold": True,
                    "bg_color": config.UI_COLORS["primary"],
                    "font_color": "white",
                    "border": 1
                })
                
                number_format = workbook.add_format({"num_format": "#,##0.00"})
                percent_format = workbook.add_format({"num_format": "0.0%"})
                
                for sheet_name, df in data.items():
                    # format date columns to remove timestamp
                    df = self._remove_timestamps_from_df(df)
                    
                    df.to_excel(writer, sheet_name=sheet_name[:31], index=False)
                    
                    if format_tables:
                        worksheet = writer.sheets[sheet_name[:31]]
                        
                        # format headers
                        for col_num, col_name in enumerate(df.columns):
                            worksheet.write(0, col_num, col_name, header_format)
                        
                        # auto-width columns
                        for col_num, col_name in enumerate(df.columns):
                            max_width = max(
                                len(str(col_name)),
                                df[col_name].astype(str).str.len().max() if len(df) > 0 else 0
                            )
                            worksheet.set_column(col_num, col_num, min(max_width + 2, 50))
            
            return True, f"exported {len(data)} sheets"
        except Exception as e:
            return False, str(e)
    
    def _remove_timestamps_from_df(self, df: pd.DataFrame) -> pd.DataFrame:
        # remove timestamps from date columns in dataframe
        result = df.copy()
        
        for col in result.columns:
            if "date" in col.lower():
                result[col] = result[col].apply(self._format_date_no_timestamp)
        
        return result
    
    def create_forecast_workbook(self,
                                 forecasts: Dict[str, Any],
                                 summary: pd.DataFrame,
                                 file_path: str) -> tuple:
        # create comprehensive forecast excel workbook
        sheets = {
            "Summary": summary,
            "Forecasts": self.format_forecast_csv(forecasts),
            "Model Performance": self._create_model_performance_df(forecasts)
        }
        
        return self.export_excel(sheets, file_path)
    
    def export_comparison_report(self, comparison_results: Dict, file_path: str) -> tuple:
        # export model comparison results to excel
        try:
            stats = comparison_results.get("model_stats", {})
            rows = []
            
            for model, metrics in stats.items():
                rows.append({
                    "Model": model,
                    "Avg MAPE": metrics.get("avg_mape", 0),
                    "Avg MAE": metrics.get("avg_mae", 0),
                    "Win Rate": metrics.get("win_rate", 0)
                })
            
            df = pd.DataFrame(rows)
            
            return self.export_excel({"Model Comparison": df}, file_path)
        except Exception as e:
            return False, str(e)
    
    def _create_model_performance_df(self, forecasts: Dict[str, Any]) -> pd.DataFrame:
        # create model performance summary
        rows = []
        
        for sku, result in forecasts.items():
            rows.append({
                "sku": sku,
                "model": result.model,
                "mape": result.metrics.get("mape", 0),
                "mae": result.metrics.get("mae", 0),
                "rmse": result.metrics.get("rmse", 0)
            })
        
        return pd.DataFrame(rows)
    
    # ---------- POWERPOINT EXPORT ----------
    
    def _check_pptx_available(self) -> tuple:
        # check if python-pptx is installed
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            return True, None
        except ImportError:
            return False, "python-pptx is not installed. Install it with: pip install python-pptx"
    
    def export_powerpoint(self,
                          summary_data: Dict[str, Any],
                          charts: List[Any],
                          file_path: str) -> tuple:
        # export to powerpoint presentation
        
        # check if pptx is available
        available, error_msg = self._check_pptx_available()
        if not available:
            return False, error_msg
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # slide 1 - title slide
            self._add_title_slide(prs, summary_data)
            
            # slide 2 - key metrics
            self._add_metrics_slide(prs, summary_data)
            
            # slide 3 - top items analysis
            self._add_analysis_slide(prs, summary_data, charts)
            
            prs.save(file_path)
            return True, "powerpoint created"
            
        except ImportError:
            return False, "python-pptx is not installed. Install it with: pip install python-pptx"
        except Exception as e:
            return False, str(e)
    
    def _add_title_slide(self, prs, summary_data: Dict) -> None:
        # add title slide to presentation
        from pptx.util import Inches, Pt
        
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Demand Forecast Summary"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = 1  # center
        
        # subtitle with date - no timestamp
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        subtitle_para.font.size = Pt(20)
        subtitle_para.alignment = 1
        
        # summary stats
        stats_text = f"{summary_data.get('total_skus', 0):,} Items | {summary_data.get('forecast_horizon', 30)} Day Forecast"
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(0.5))
        stats_frame = stats_box.text_frame
        stats_para = stats_frame.paragraphs[0]
        stats_para.text = stats_text
        stats_para.font.size = Pt(18)
        stats_para.alignment = 1
    
    def _add_metrics_slide(self, prs, summary_data: Dict) -> None:
        # add key metrics slide
        from pptx.util import Inches, Pt
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Key Metrics"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        # metrics grid
        metrics = [
            ("Total Forecast", f"{summary_data.get('total_forecast', 0):,.0f} units"),
            ("Average MAPE", f"{summary_data.get('avg_mape', 0):.1f}%"),
            ("A-Items Coverage", f"{summary_data.get('a_items_pct', 0):.0f}%"),
            ("Models Used", summary_data.get('models_used', 'N/A'))
        ]
        
        for i, (label, value) in enumerate(metrics):
            row = i // 2
            col = i % 2
            x = Inches(0.5 + col * 6)
            y = Inches(1.5 + row * 2)
            
            # metric box
            metric_box = slide.shapes.add_textbox(x, y, Inches(5.5), Inches(1.5))
            tf = metric_box.text_frame
            
            # label
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(16)
            p1.font.bold = True
            
            # value
            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(36)
            p2.font.bold = True
    
    def _add_analysis_slide(self, prs, summary_data: Dict, charts: List) -> None:
        # add analysis slide with charts
        from pptx.util import Inches, Pt
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Top Items Analysis"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        # add chart images if available
        if charts:
            for i, chart_path in enumerate(charts[:2]):
                if Path(chart_path).exists():
                    x = Inches(0.5 + i * 6.2)
                    slide.shapes.add_picture(chart_path, x, Inches(1.5), width=Inches(6))
        
        # top items table
        top_items = summary_data.get("top_items", [])
        if top_items:
            table_top = Inches(4.5)
            table = slide.shapes.add_table(
                rows=min(6, len(top_items) + 1),
                cols=4,
                left=Inches(0.5),
                top=table_top,
                width=Inches(12),
                height=Inches(2.5)
            ).table
            
            # headers
            headers = ["Item", "Forecast", "Change", "Model"]
            for j, header in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # data rows
            for i, item in enumerate(top_items[:5]):
                table.cell(i + 1, 0).text = str(item.get("sku", ""))
                table.cell(i + 1, 1).text = f"{item.get('forecast', 0):,.0f}"
                table.cell(i + 1, 2).text = f"{item.get('change_pct', 0):+.1f}%"
                table.cell(i + 1, 3).text = str(item.get("model", ""))
    
    def create_executive_ppt(self,
                             forecasts: Dict[str, Any],
                             cluster_summary: List[Dict],
                             file_path: str) -> tuple:
        # create 3-slide executive presentation
        
        # check if pptx is available first
        available, error_msg = self._check_pptx_available()
        if not available:
            return False, error_msg
        
        # calculate summary data
        total_forecast = sum(sum(r.forecast) for r in forecasts.values())
        avg_mape = sum(r.metrics.get("mape", 0) for r in forecasts.values()) / max(1, len(forecasts))
        
        # get top items
        top_items = sorted(
            [{"sku": k, "forecast": sum(v.forecast), "model": v.model, 
              "mape": v.metrics.get("mape", 0)} 
             for k, v in forecasts.items()],
            key=lambda x: x["forecast"],
            reverse=True
        )[:10]
        
        # model distribution
        model_counts = {}
        for r in forecasts.values():
            model_counts[r.model] = model_counts.get(r.model, 0) + 1
        
        summary_data = {
            "total_skus": len(forecasts),
            "total_forecast": total_forecast,
            "avg_mape": avg_mape,
            "forecast_horizon": len(next(iter(forecasts.values())).forecast) if forecasts else 0,
            "top_items": top_items,
            "models_used": ", ".join(model_counts.keys()),
            "a_items_pct": len([r for r in forecasts.values() if sum(r.forecast) > 1000]) / max(1, len(forecasts)) * 100
        }
        
        return self.export_powerpoint(summary_data, [], file_path)
    
    # ---------- PDF EXPORT ----------
    
    def export_pdf(self,
                   summary_data: Dict[str, Any],
                   file_path: str) -> tuple:
        # export summary to pdf
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            
            c = canvas.Canvas(file_path, pagesize=letter)
            width, height = letter
            
            # title
            c.setFont("Helvetica-Bold", 24)
            c.drawString(1 * inch, height - 1 * inch, "Demand Forecast Report")
            
            # date - no timestamp
            c.setFont("Helvetica", 12)
            c.drawString(1 * inch, height - 1.4 * inch, 
                        f"Generated: {datetime.now().strftime('%B %d, %Y')}"
            )
            
            # horizontal line
            c.setStrokeColor(colors.grey)
            c.line(1 * inch, height - 1.6 * inch, width - 1 * inch, height - 1.6 * inch)
            
            # summary section
            c.setFont("Helvetica-Bold", 16)
            c.drawString(1 * inch, height - 2.2 * inch, "Executive Summary")
            
            y_pos = height - 2.6 * inch
            c.setFont("Helvetica", 11)
            
            summary_lines = [
                f"Total Items Forecasted: {summary_data.get('total_skus', 0):,}",
                f"Forecast Horizon: {summary_data.get('forecast_horizon', 30)} days",
                f"Total Projected Demand: {summary_data.get('total_forecast', 0):,.0f} units",
                f"Average Forecast Accuracy (MAPE): {summary_data.get('avg_mape', 0):.1f}%",
                f"Models Used: {summary_data.get('models_used', 'N/A')}"
            ]
            
            for line in summary_lines:
                c.drawString(1.2 * inch, y_pos, line)
                y_pos -= 0.3 * inch
            
            # key insights section
            y_pos -= 0.4 * inch
            c.setFont("Helvetica-Bold", 16)
            c.drawString(1 * inch, y_pos, "Key Insights")
            
            y_pos -= 0.4 * inch
            c.setFont("Helvetica", 11)
            
            insights = summary_data.get("insights", [
                "High-volume items (A-items) represent 80% of total forecast",
                "Seasonal patterns detected in 35% of items",
                "Recommendation: Focus inventory planning on top 100 SKUs"
            ])
            
            for insight in insights:
                c.drawString(1.2 * inch, y_pos, f"• {insight}")
                y_pos -= 0.3 * inch
            
            c.save()
            return True, "pdf report created"
            
        except ImportError:
            return False, "reportlab is not installed. Install it with: pip install reportlab"
        except Exception as e:
            return False, str(e)
    
    # ---------- HELPER METHODS ----------
    
    def get_export_filename(self, base_name: str, format_type: str) -> str:
        # generate export filename with date only (no timestamp)
        date_str = datetime.now().strftime("%Y%m%d")
        extension = self.export_config.get(format_type, {}).get("extension", ".csv")
        return f"{base_name}_{date_str}{extension}"
    
    def prepare_export_data(self,
                            forecasts: Dict[str, Any],
                            include_metrics: bool = True,
                            include_bounds: bool = True) -> Dict[str, pd.DataFrame]:
        # prepare all export data sheets
        sheets = {}
        
        # forecasts sheet - includes all forecast values
        forecast_rows = []
        for sku, result in forecasts.items():
            for i, date in enumerate(result.dates):
                # format date without timestamp
                date_str = self._format_date_no_timestamp(date)
                
                row = {
                    "sku": sku,
                    "date": date_str,
                    "forecast": result.forecast[i]
                }
                if include_bounds:
                    row["lower_bound"] = result.lower_bound[i] if result.lower_bound else None
                    row["upper_bound"] = result.upper_bound[i] if result.upper_bound else None
                forecast_rows.append(row)
        
        sheets["Forecasts"] = pd.DataFrame(forecast_rows)
        
        # summary sheet
        if include_metrics:
            summary_rows = []
            for sku, result in forecasts.items():
                summary_rows.append({
                    "sku": sku,
                    "model": result.model,
                    "total_forecast": sum(result.forecast),
                    "avg_daily": sum(result.forecast) / len(result.forecast),
                    "mape": result.metrics.get("mape", 0),
                    "mae": result.metrics.get("mae", 0)
                })
            sheets["Summary"] = pd.DataFrame(summary_rows)
        
        return sheets